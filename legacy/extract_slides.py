import sys
from pptx import Presentation

def extract_slides(input_file, slide_numbers, output_file):
    """
    Извлекает указанные слайды из PPTX и сохраняет их в новый файл.

    Параметры:
        input_file (str): путь к исходному файлу.
        slide_numbers (list): номера слайдов для сохранения (отсчёт с 1).
        output_file (str): путь для результата.
    """
    try:
        prs = Presentation(input_file)
    except Exception as e:
        print(f"Ошибка при открытии файла: {e}")
        sys.exit(1)

    total_slides = len(prs.slides)
    # Преобразуем номера в индексы и фильтруем допустимые
    indices = sorted({num - 1 for num in slide_numbers if 1 <= num <= total_slides})

    if not indices:
        print("Нет корректных номеров слайдов.")
        sys.exit(1)

    # Получаем идентификаторы слайдов, которые нужно оставить
    keep_ids = [prs.slides[i].slide_id for i in indices]

    # Работаем напрямую со списком ссылок на слайды
    sld_id_lst = prs.slides._sldIdLst
    # Итерируемся по копии списка, чтобы безопасно удалять элементы
    for sld_id in list(sld_id_lst):
        if sld_id.id not in keep_ids:
            sld_id_lst.remove(sld_id)

    prs.save(output_file)
    print(f"Готово. Сохранено {len(indices)} слайдов в '{output_file}'.")

if __name__ == "__main__":
    # Пример вызова
    extract_slides("docs/111/input.pptx", [4], "output1.pptx")