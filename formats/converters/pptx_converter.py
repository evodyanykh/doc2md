"""
Конвертер PPTX → Markdown.

Ключевая сложность PPTX — это отсутствие явного порядка чтения.
Слайд содержит набор фигур (shapes) в произвольном порядке,
расположенных по координатам (left, top).

Стратегия порядка чтения:
  1. Заголовок слайда (placeholder type TITLE/CENTER_TITLE) — всегда первый
  2. Остальные фигуры сортируются по строкам (top) с допуском на выравнивание,
     внутри строки — слева направо (left).
     Это обрабатывает как вертикальный макет, так и колонки.

Сохраняемое форматирование:
  - Заголовок слайда → # Заголовок
  - Подзаголовки → ## Подзаголовок
  - Жирный текст → **текст**
  - Курсив → *текст*
  - Маркированные списки с уровнями вложенности → - / двойные отступы
  - Таблицы → Markdown-таблица
  - Изображения → плейсхолдер или OCR

Библиотека: python-pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

from formats.config import IMAGE_PLACEHOLDER, TABLE_WARNING
from formats.converters.base import BaseConverter, ConversionResult
from formats.utils.ocr import ocr_image_bytes, is_tesseract_available


# Допуск по вертикали для определения одной «строки» (в EMU — единицы pptx)
# 914400 EMU = 1 дюйм. 500000 ≈ 0.5 дюйма — достаточно для группировки в строку.
ROW_TOLERANCE_EMU = 500_000


class PptxConverter(BaseConverter):

    def convert(self, file_path: Path) -> ConversionResult:
        result = ConversionResult()
        prs = Presentation(str(file_path))

        all_slides: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_lines = self._convert_slide(slide, slide_num, result)
            all_slides.append("\n".join(slide_lines))

        # Слайды разделяем горизонтальной линией
        result.content = self._clean_text("\n\n---\n\n".join(all_slides))
        return result

    # ─── Слайд ────────────────────────────────────────────────────────────────

    def _convert_slide(self, slide, slide_num: int, result: ConversionResult) -> list[str]:
        """Конвертирует один слайд в список Markdown-строк."""
        lines: list[str] = [f"## Слайд {slide_num}"]
        lines.append("")

        # Разделяем заголовок и остальные фигуры
        title_shape = self._find_title(slide)
        other_shapes = [s for s in slide.shapes if s != title_shape]

        # Заголовок слайда
        if title_shape and title_shape.has_text_frame:
            title_text = self._get_plain_text(title_shape)
            if title_text:
                lines.append(f"# {title_text}")
                lines.append("")

        # Сортируем фигуры по строкам (top), внутри строки — по left
        sorted_shapes = self._sort_shapes_by_reading_order(other_shapes)

        for shape in sorted_shapes:
            shape_lines = self._convert_shape(shape, result)
            if shape_lines:
                lines.extend(shape_lines)
                lines.append("")

        return lines

    # ─── Сортировка фигур (логика порядка чтения) ─────────────────────────────

    def _sort_shapes_by_reading_order(self, shapes) -> list:
        """
        Сортирует фигуры в порядке чтения: сверху вниз, слева направо.

        Алгоритм:
          1. Группируем фигуры в «строки» по вертикальной позиции (top) с допуском ROW_TOLERANCE_EMU.
             Фигуры в одной строке — это колонки.
          2. Строки сортируем по минимальному top.
          3. Внутри каждой строки сортируем по left.
        """
        # Фильтруем фигуры без позиции (например, группы без координат)
        valid_shapes = [s for s in shapes if hasattr(s, "top") and s.top is not None]

        if not valid_shapes:
            return []

        # Группировка в строки
        rows: list[list] = []
        for shape in sorted(valid_shapes, key=lambda s: s.top):
            placed = False
            for row in rows:
                # Сравниваем с минимальным top строки
                row_top = min(s.top for s in row)
                if abs(shape.top - row_top) <= ROW_TOLERANCE_EMU:
                    row.append(shape)
                    placed = True
                    break
            if not placed:
                rows.append([shape])

        # Сортируем строки по top и фигуры внутри строки по left
        rows.sort(key=lambda row: min(s.top for s in row))
        result: list = []
        for row in rows:
            result.extend(sorted(row, key=lambda s: s.left if s.left is not None else 0))

        return result

    # ─── Фигуры ───────────────────────────────────────────────────────────────

    def _convert_shape(self, shape, result: ConversionResult) -> list[str]:
        """Определяет тип фигуры и вызывает нужный метод конвертации."""

        # Таблица
        if shape.has_table:
            lines = self._convert_table(shape.table, result)
            return lines

        # Текстовый фрейм
        if shape.has_text_frame:
            return self._convert_text_frame(shape.text_frame, shape, result)

        # Изображение (Picture)
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE = 13
            return self._convert_image(shape, result)

        return []

    # ─── Текст ────────────────────────────────────────────────────────────────

    def _convert_text_frame(self, text_frame, shape, result: ConversionResult) -> list[str]:
        """
        Конвертирует TextFrame (набор параграфов внутри фигуры).
        Определяет является ли фигура подзаголовком.
        """
        lines: list[str] = []
        is_subtitle = self._is_subtitle(shape)

        for para in text_frame.paragraphs:
            md_line = self._convert_paragraph(para, is_subtitle)
            if md_line is not None:
                lines.append(md_line)

        return lines

    def _convert_paragraph(self, para, is_subtitle: bool = False) -> str | None:
        """
        Конвертирует параграф PPTX в Markdown-строку.

        para.level — уровень вложенности списка (0 = верхний уровень).
        """
        text = self._convert_runs(para)

        # Фильтруем пустой текст и строки-артефакты вида "- -", "**  **"
        if not "".join(text.split()).replace("-", "").replace("*", ""):
            return ""

        # Определяем отступ по уровню вложенности
        indent = "  " * para.level

        # Подзаголовок (второй placeholder) → ##
        if is_subtitle and para.level == 0:
            return f"## {text}"

        # Маркированный список: level > 0 или текст начинается с маркера
        if para.level > 0:
            return f"{indent}- {text}"

        # Определяем по размеру шрифта: крупный текст = заголовочный
        if self._is_large_text(para):
            return f"### {text}"

        # Обычный текст или пункт списка верхнего уровня
        # Проверяем признаки маркированного списка через XML
        if self._has_bullet(para):
            return f"- {text}"

        return text

    def _convert_runs(self, para) -> str:
        """Собирает текст из runs параграфа с форматированием bold/italic."""
        parts: list[str] = []

        for run in para.runs:
            text = run.text
            if not text:
                continue

            bold = run.font.bold
            italic = run.font.italic

            # Применяем форматирование только к непустому тексту —
            # иначе получаются артефакты **** или **  **
            if text.strip():
                if bold and italic:
                    text = f"***{text}***"
                elif bold:
                    text = f"**{text}**"
                elif italic:
                    text = f"*{text}*"

            parts.append(text)

        return "".join(parts)

    def _is_large_text(self, para) -> bool:
        """
        Эвристика: если средний размер шрифта в параграфе >= 24pt —
        считаем его заголовочным.
        """
        sizes = [
            run.font.size.pt
            for run in para.runs
            if run.font.size is not None
        ]
        if not sizes:
            return False
        return (sum(sizes) / len(sizes)) >= 24

    def _has_bullet(self, para) -> bool:
        """
        Проверяет наличие буллета через XML.
        <a:buChar> или <a:buAutoNum> означают маркированный/нумерованный список.
        """
        pPr = para._p.find("{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
        if pPr is None:
            return False
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        return (
            pPr.find(f"{{{ns}}}buChar") is not None
            or pPr.find(f"{{{ns}}}buAutoNum") is not None
            or pPr.find(f"{{{ns}}}buClr") is not None
        )

    # ─── Заголовок слайда ─────────────────────────────────────────────────────

    def _find_title(self, slide):
        """
        Ищет заголовок слайда среди placeholders.
        Типы: TITLE (1), CENTER_TITLE (3).
        """
        for shape in slide.placeholders:
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                return shape
        return None

    def _is_subtitle(self, shape) -> bool:
        """Проверяет, является ли фигура подзаголовком (SUBTITLE placeholder)."""
        if not hasattr(shape, "placeholder_format") or shape.placeholder_format is None:
            return False
        return shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE

    def _get_plain_text(self, shape) -> str:
        """Извлекает весь текст из фигуры без форматирования."""
        if shape.has_text_frame:
            return "\n".join(
                para.text for para in shape.text_frame.paragraphs if para.text.strip()
            )
        return ""

    # ─── Изображения ──────────────────────────────────────────────────────────

    def _convert_image(self, shape, result: ConversionResult) -> list[str]:
        """
        Обрабатывает изображение на слайде.
        Пытается извлечь текст через OCR.
        """
        tesseract_ok = is_tesseract_available()

        try:
            image_bytes = shape.image.blob

            if tesseract_ok:
                ocr_text = ocr_image_bytes(image_bytes)
                if ocr_text and IMAGE_PLACEHOLDER not in ocr_text:
                    return [f"\n> *[Текст из изображения]:* {ocr_text}\n"]
                else:
                    result.add_warning(f"Изображение без текста на слайде")
                    return [IMAGE_PLACEHOLDER]
            else:
                result.add_warning("Tesseract недоступен — изображение пропущено")
                return [IMAGE_PLACEHOLDER]

        except Exception as e:
            result.add_warning(f"Не удалось обработать изображение: {e}")
            return [IMAGE_PLACEHOLDER]

    # ─── Таблицы ──────────────────────────────────────────────────────────────

    def _convert_table(self, table, result: ConversionResult) -> list[str]:
        """Конвертирует таблицу PPTX в Markdown."""
        rows: list[list[str]] = []

        for row in table.rows:
            cells: list[str] = []
            for cell in row.cells:
                cell_text = cell.text.replace("\n", " ").strip()
                cells.append(cell_text)
            rows.append(cells)

        if not rows:
            return []

        col_count = max(len(row) for row in rows)
        for row in rows:
            while len(row) < col_count:
                row.append("")

        md_lines: list[str] = []
        md_lines.append("| " + " | ".join(rows[0]) + " |")
        md_lines.append("| " + " | ".join(["---"] * col_count) + " |")
        for row in rows[1:]:
            md_lines.append("| " + " | ".join(row) + " |")

        if col_count > 2:
            result.add_warning(f"Таблица {len(rows)}×{col_count}: {TABLE_WARNING}")

        return md_lines
